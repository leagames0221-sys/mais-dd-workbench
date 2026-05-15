"""Synthetic VDR data generator (T2 Week 1、 ADR-100/101/102/103 SSoT 実装).

Run: python -m src.data_gen.generate_synthetic_vdr

Output (data/vdr_synthetic/<DDP-id>/<doc_kind>/<filename>):
  - 5 DDP 案件 (M&A 試作対象企業)
  - 各 DDP × 8 docs (PDF 2 + DOCX 2 + XLSX 2 + PPTX 2) = 計 40 docs total
  - 中堅日本企業 fit pattern (同族経営 / 名義株 / オーナー私的経費) を literal inject

ADR-006 PII boundary 順守:
  - 合成 data のみ生成、 vault 想定 PII field (担当者連絡先) は raw text 内 inject 後、
    Docling ingestion 時の redact layer で literal 抹消される設計
  - 本 generator は raw + redacted の両 layer test 用途も literal cover
"""
from __future__ import annotations

import os
import random
import secrets
import sys

# Windows cp932 default encoding 防御 (T1 security_intercept.py と同 pattern、 doctrine: analogical-recall)
# ¥ 記号や 日本語 を print する際に literal 必要、 fix なしだと cp932 encoding error abort
for _stream_name in ("stdout", "stderr"):
    _stream = getattr(sys, _stream_name, None)
    if _stream is not None and getattr(_stream, "encoding", "").lower() != "utf-8":
        try:
            _stream.reconfigure(encoding="utf-8", errors="replace")
        except Exception:
            pass
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Callable

from dotenv import load_dotenv
from faker import Faker
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.platypus import (
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)
from reportlab.lib import colors

# python-docx は名前空間衝突を避けるため import 形式に注意
import docx as python_docx_module

from src.data_gen.jp_patterns import (
    JPPatternSeed,
    industry_median_compensation_jpy,
    industry_revenue_jpy,
    make_seed,
)

load_dotenv()

# ===== const =====
SEED = int(os.environ.get("SYNTHETIC_SEED", "20260512"))
DDP_COUNT = int(os.environ.get("SYNTHETIC_DDP_COUNT", "5"))
OUTPUT_ROOT = Path(os.environ.get("VDR_OUTPUT_ROOT", "./data/vdr_synthetic"))

# T1 INDUSTRIES から literal reuse (mais-deal-matching/src/data_gen)
INDUSTRIES = [
    "出版", "映像制作", "広告", "印刷", "アパレル",
    "食品製造", "外食", "小売", "卸売", "物流",
    "建設", "不動産", "ホテル", "観光", "教育",
    "医療", "介護", "農業", "IT", "コンサルティング",
]

PREFECTURES = [
    "東京都", "大阪府", "京都府", "神奈川県", "愛知県", "兵庫県",
    "福岡県", "北海道", "千葉県", "埼玉県", "長野県", "静岡県",
]

# reportlab で日本語表示 (CID font、 reportlab 公式同梱)
pdfmetrics.registerFont(UnicodeCIDFont("HeiseiKakuGo-W5"))
JP_FONT = "HeiseiKakuGo-W5"


# ===== DDProject schema =====


@dataclass
class DDProject:
    """1 M&A DD 案件 (DDP-XXXXXX)."""

    ddp_id: str
    company_name: str
    industry: str
    revenue_jpy: int
    founder_name: str  # 創業者氏名 (姓 = 同族経営時 family_surname)
    founder_age: int
    location_pref: str
    employee_count: int
    contact_email: str  # PII (vault 行き想定)
    contact_phone: str  # PII (vault 行き想定)
    jp_seed: JPPatternSeed
    directors: list[str] = field(default_factory=list)  # 役員氏名 list
    shareholders: list[tuple[str, float]] = field(default_factory=list)
    # shareholders = [(name, share_ratio_percent), ...]


def make_ddp(idx: int, rng: random.Random, faker: Faker) -> DDProject:
    """1 DDP 案件 を literal 合成."""
    industry = rng.choice(INDUSTRIES)
    revenue = industry_revenue_jpy(industry, rng)
    jp_seed = make_seed(industry, rng)

    # founder 姓 = jp_seed.family_surname が literal あればそれ、 そうでなければ Faker
    if jp_seed.family_surname:
        founder_surname = jp_seed.family_surname
    else:
        founder_surname = faker.last_name()
    founder_given = rng.choice(["太郎", "正一", "幸夫", "茂", "和男", "健司", "誠"])
    founder_name = f"{founder_surname} {founder_given}"

    # 役員: 同族経営 = family_members 多数、 そうでなければ Faker 由来
    directors: list[str] = [founder_name]
    if jp_seed.family_governance:
        directors.extend(jp_seed.family_members[:3])
    else:
        for _ in range(3):
            directors.append(f"{faker.last_name()} {faker.first_name()}")

    # 株主: 同族経営 = family_members で 60-85% 占有、 残りは外部
    shareholders: list[tuple[str, float]] = []
    if jp_seed.family_governance:
        family_total_ratio = rng.uniform(60.0, 85.0)
        members_for_share = [founder_name] + jp_seed.family_members[:4]
        per = family_total_ratio / len(members_for_share)
        for name in members_for_share:
            shareholders.append((name, round(per, 2)))
        remaining = 100.0 - family_total_ratio
        # 外部株主 2-3 名
        ext_count = rng.randint(2, 3)
        per_ext = remaining / ext_count
        for _ in range(ext_count):
            shareholders.append((f"{faker.last_name()} {faker.first_name()}", round(per_ext, 2)))
    else:
        # 分散型: 創業者 35-50% + 外部 5-7 名
        founder_share = rng.uniform(35.0, 50.0)
        shareholders.append((founder_name, round(founder_share, 2)))
        remaining = 100.0 - founder_share
        ext_count = rng.randint(5, 7)
        per_ext = remaining / ext_count
        for _ in range(ext_count):
            shareholders.append((f"{faker.last_name()} {faker.first_name()}", round(per_ext, 2)))

    company_suffix = rng.choice(["株式会社", "有限会社"])
    # Faker.company() は "合同会社XX" / "株式会社XX" / "有限会社XX" / "XX" を return、
    # 全 prefix を literal 除去してから 自分の prefix を 1 つだけ付与 (suffix duplication 防御)
    raw_company = faker.company()
    for prefix in ("合同会社", "株式会社", "有限会社", "合資会社", "合名会社"):
        raw_company = raw_company.replace(prefix, "")
    company_name = f"{company_suffix}{raw_company.strip()}"

    return DDProject(
        ddp_id=f"DDP-{secrets.token_urlsafe(8)}",
        company_name=company_name,
        industry=industry,
        revenue_jpy=revenue,
        founder_name=founder_name,
        founder_age=rng.randint(58, 78),
        location_pref=rng.choice(PREFECTURES),
        employee_count=rng.randint(30, 250),
        contact_email=faker.email(),
        contact_phone=faker.phone_number(),
        jp_seed=jp_seed,
        directors=directors,
        shareholders=shareholders,
    )


# ===== PDF generators (reportlab) =====


def _pdf_styles() -> dict:
    """日本語対応 paragraph styles."""
    base = getSampleStyleSheet()
    base["Title"].fontName = JP_FONT
    base["Title"].fontSize = 16
    base["Normal"].fontName = JP_FONT
    base["Normal"].fontSize = 10
    base["Heading1"].fontName = JP_FONT
    base["Heading2"].fontName = JP_FONT
    return base


def generate_pdf_tax_return(ddp: DDProject, out_path: Path, rng: random.Random) -> None:
    """法人税申告書 (簡易) sample PDF を literal 生成."""
    out_path.parent.mkdir(parents=True, exist_ok=True)
    doc = SimpleDocTemplate(str(out_path), pagesize=A4, topMargin=20 * mm, bottomMargin=20 * mm)
    styles = _pdf_styles()
    story: list = []

    story.append(Paragraph(f"法人税申告書 (合成 sample) — {ddp.company_name}", styles["Title"]))
    story.append(Spacer(1, 8 * mm))
    story.append(Paragraph(f"事業年度: 2024-04-01 ~ 2025-03-31", styles["Normal"]))
    story.append(Paragraph(f"所在地: {ddp.location_pref}", styles["Normal"]))
    story.append(Paragraph(f"業種: {ddp.industry}", styles["Normal"]))
    story.append(Spacer(1, 6 * mm))

    # 損益計算書 簡易 table
    revenue = ddp.revenue_jpy
    cogs = int(revenue * rng.uniform(0.55, 0.70))
    gross = revenue - cogs
    median_comp = industry_median_compensation_jpy(ddp.industry)
    officer_comp = int(median_comp * ddp.jp_seed.owner_compensation_multiplier)
    sga = int(revenue * rng.uniform(0.15, 0.22)) + officer_comp
    operating = gross - sga
    pretax = operating - rng.randint(5_000_000, 30_000_000)
    tax = max(0, int(pretax * 0.30))
    net = pretax - tax

    pl_data = [
        ["項目", "金額 (円)"],
        ["売上高", f"{revenue:,}"],
        ["売上原価", f"{cogs:,}"],
        ["売上総利益", f"{gross:,}"],
        ["役員報酬 (合計)", f"{officer_comp:,}"],
        ["販管費 (合計)", f"{sga:,}"],
        ["営業利益", f"{operating:,}"],
        ["税引前当期純利益", f"{pretax:,}"],
        ["法人税等", f"{tax:,}"],
        ["当期純利益", f"{net:,}"],
    ]
    if ddp.jp_seed.owner_private_expense:
        pl_data.append(["役員貸付金 (残高)", f"{ddp.jp_seed.owner_loan_amount_jpy:,}"])
        pl_data.append(
            ["備考", f"役員報酬は業界中央値 (¥{median_comp:,}) の {ddp.jp_seed.owner_compensation_multiplier} 倍"]
        )
    table = Table(pl_data, colWidths=[80 * mm, 80 * mm])
    table.setStyle(
        TableStyle(
            [
                ("FONT", (0, 0), (-1, -1), JP_FONT, 9),
                ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                ("GRID", (0, 0), (-1, -1), 0.3, colors.grey),
                ("ALIGN", (1, 1), (1, -1), "RIGHT"),
            ]
        )
    )
    story.append(table)
    story.append(Spacer(1, 4 * mm))
    story.append(Paragraph(f"代表者: {ddp.founder_name}", styles["Normal"]))
    story.append(Paragraph(f"連絡先: {ddp.contact_email} / {ddp.contact_phone}", styles["Normal"]))
    doc.build(story)


def generate_pdf_financial_statement(ddp: DDProject, out_path: Path, rng: random.Random) -> None:
    """財務諸表 (BS 簡易) sample PDF を literal 生成."""
    out_path.parent.mkdir(parents=True, exist_ok=True)
    doc = SimpleDocTemplate(str(out_path), pagesize=A4, topMargin=20 * mm, bottomMargin=20 * mm)
    styles = _pdf_styles()
    story: list = []

    story.append(Paragraph(f"貸借対照表 (合成 sample) — {ddp.company_name}", styles["Title"]))
    story.append(Spacer(1, 8 * mm))
    story.append(Paragraph("基準日: 2025-03-31", styles["Normal"]))
    story.append(Spacer(1, 6 * mm))

    cash = int(ddp.revenue_jpy * rng.uniform(0.10, 0.20))
    ar = int(ddp.revenue_jpy * rng.uniform(0.12, 0.18))
    inventory = int(ddp.revenue_jpy * rng.uniform(0.08, 0.15))
    ppe = int(ddp.revenue_jpy * rng.uniform(0.20, 0.35))
    total_assets = cash + ar + inventory + ppe

    bs_data = [
        ["資産項目", "金額 (円)", "負債・純資産", "金額 (円)"],
        ["現金及び預金", f"{cash:,}", "買掛金", f"{int(total_assets * 0.18):,}"],
        ["売掛金", f"{ar:,}", "短期借入金", f"{int(total_assets * 0.12):,}"],
        ["棚卸資産", f"{inventory:,}", "長期借入金", f"{int(total_assets * 0.20):,}"],
        ["有形固定資産", f"{ppe:,}", "資本金", f"{int(total_assets * 0.15):,}"],
    ]
    if ddp.jp_seed.owner_private_expense:
        bs_data.append(["役員貸付金", f"{ddp.jp_seed.owner_loan_amount_jpy:,}", "利益剰余金", f"{int(total_assets * 0.35):,}"])
    else:
        bs_data.append(["その他資産", f"{int(total_assets * 0.05):,}", "利益剰余金", f"{int(total_assets * 0.35):,}"])
    bs_data.append(["合計", f"{total_assets:,}", "合計", f"{total_assets:,}"])

    table = Table(bs_data, colWidths=[45 * mm, 40 * mm, 45 * mm, 40 * mm])
    table.setStyle(
        TableStyle(
            [
                ("FONT", (0, 0), (-1, -1), JP_FONT, 9),
                ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                ("GRID", (0, 0), (-1, -1), 0.3, colors.grey),
                ("ALIGN", (1, 1), (1, -1), "RIGHT"),
                ("ALIGN", (3, 1), (3, -1), "RIGHT"),
            ]
        )
    )
    story.append(table)
    story.append(Spacer(1, 4 * mm))
    story.append(Paragraph(f"監査人連絡先: {ddp.contact_email}", styles["Normal"]))
    doc.build(story)


# ===== DOCX generators (python-docx) =====


def generate_docx_minutes(ddp: DDProject, out_path: Path, rng: random.Random) -> None:
    """取締役会議事録 sample DOCX を literal 生成 (同族経営時 = 議論欠如、 通常時 = 健全議論)."""
    out_path.parent.mkdir(parents=True, exist_ok=True)
    document = python_docx_module.Document()
    document.add_heading(f"取締役会議事録 — {ddp.company_name}", level=1)
    document.add_paragraph(f"開催日: 2025-{rng.randint(1, 12):02d}-{rng.randint(1, 28):02d}")
    document.add_paragraph(f"出席役員: {', '.join(ddp.directors)}")
    document.add_paragraph(f"場所: 本社会議室 ({ddp.location_pref})")

    document.add_heading("議題 1: 来期事業計画", level=2)
    if ddp.jp_seed.family_governance:
        # 同族経営: 議論欠如 pattern (代表が一方的に説明、 異論 ZERO)
        document.add_paragraph(f"代表取締役 {ddp.founder_name} より来期計画の説明あり。 全員一致で承認。")
        document.add_paragraph(f"(議論記録 なし)")
    else:
        # 健全議論
        document.add_paragraph(
            f"代表取締役 {ddp.founder_name} より来期計画の説明後、 各取締役より以下の質問:"
        )
        document.add_paragraph(
            f"・{ddp.directors[1]}: 新規市場参入の COGS 影響を再試算すべきでは"
        )
        document.add_paragraph(
            f"・{ddp.directors[2]}: 投資回収期間が 5 年想定だが、 為替 risk を反映した sensitivity 分析を依頼"
        )
        document.add_paragraph(f"修正提案を反映の上、 次回再審議とする。")

    document.add_heading("議題 2: 役員報酬", level=2)
    median = industry_median_compensation_jpy(ddp.industry)
    officer_comp = int(median * ddp.jp_seed.owner_compensation_multiplier)
    document.add_paragraph(f"代表取締役 {ddp.founder_name} の年間役員報酬を ¥{officer_comp:,} に決定。")
    if ddp.jp_seed.owner_private_expense:
        document.add_paragraph(
            f"(参考: 業界中央値は ¥{median:,})、 比率 {ddp.jp_seed.owner_compensation_multiplier} 倍"
        )

    if ddp.jp_seed.nominee_shareholder:
        document.add_heading("議題 3: 株主名義変更承認", level=2)
        for frm, to, dt in ddp.jp_seed.nominee_chain:
            document.add_paragraph(f"・{dt}: {frm} → {to} 名義変更 (家族間譲渡)")

    document.add_paragraph(f"連絡先: {ddp.contact_email}")
    document.save(str(out_path))


def generate_docx_regulations(ddp: DDProject, out_path: Path, rng: random.Random) -> None:
    """株式譲渡制限規程 + Change of Control 条項 sample DOCX を literal 生成."""
    out_path.parent.mkdir(parents=True, exist_ok=True)
    document = python_docx_module.Document()
    document.add_heading(f"株式譲渡制限規程 — {ddp.company_name}", level=1)
    document.add_paragraph(f"制定日: 2018-04-01、 改訂日: 2024-{rng.randint(1, 12):02d}-01")

    document.add_heading("第 1 条 (目的)", level=2)
    document.add_paragraph(f"本規程は {ddp.company_name} (以下、 「当社」) の株式譲渡を制限することにより、 当社の経営基盤の安定を図ることを目的とする。")

    document.add_heading("第 2 条 (譲渡の承認)", level=2)
    document.add_paragraph("株主が所有する株式の全部又は一部を譲渡しようとするときは、 取締役会の承認を要するものとする。")

    document.add_heading("第 3 条 (Change of Control 条項)", level=2)
    document.add_paragraph(
        "発行済株式総数の 50% を超える持分異動が生じた場合、 当社の取引先及び金融機関に対し速やかに通知するものとする。"
        " また、 当社が締結する主要取引契約のうち、 別表 A に列挙する契約は Change of Control 条項を含むため、"
        " 経営権異動の前に相手方の書面同意を取得することを要する。"
    )

    document.add_heading("別表 A: Change of Control 条項を含む主要契約", level=2)
    document.add_paragraph(f"・主要取引先 (年間取引額 ¥{int(ddp.revenue_jpy * 0.15):,} 以上) との基本契約 3 件")
    document.add_paragraph("・主要金融機関 (○○銀行) との融資契約")
    document.add_paragraph("・販売代理店契約 (海外)")

    if ddp.jp_seed.nominee_shareholder:
        document.add_heading("補足: 過去の名義変更履歴 (家族間)", level=2)
        for frm, to, dt in ddp.jp_seed.nominee_chain:
            document.add_paragraph(f"・{dt}: 株主 {frm} → {to} に名義変更")

    document.save(str(out_path))


# ===== XLSX generators (openpyxl) =====


def generate_xlsx_trial_balance(ddp: DDProject, out_path: Path, rng: random.Random) -> None:
    """月次試算表 sample XLSX を literal 生成."""
    out_path.parent.mkdir(parents=True, exist_ok=True)
    wb = Workbook()
    ws = wb.active
    ws.title = "試算表"

    ws["A1"] = f"月次試算表 — {ddp.company_name}"
    ws["A1"].font = Font(bold=True, size=14)
    ws["A2"] = "基準月: 2025-03"
    ws["A4"] = "勘定科目"
    ws["B4"] = "借方"
    ws["C4"] = "貸方"
    header_fill = PatternFill("solid", fgColor="DDDDDD")
    for col in ["A4", "B4", "C4"]:
        ws[col].font = Font(bold=True)
        ws[col].fill = header_fill

    rows = [
        ("現金及び預金", int(ddp.revenue_jpy * 0.15 / 12), 0),
        ("売掛金", int(ddp.revenue_jpy * 0.18 / 12), 0),
        ("買掛金", 0, int(ddp.revenue_jpy * 0.12 / 12)),
        ("売上高", 0, int(ddp.revenue_jpy / 12)),
        ("売上原価", int(ddp.revenue_jpy * 0.6 / 12), 0),
        ("役員報酬",
         int(industry_median_compensation_jpy(ddp.industry) * ddp.jp_seed.owner_compensation_multiplier / 12),
         0),
        ("人件費", int(ddp.revenue_jpy * 0.18 / 12), 0),
        ("地代家賃", int(ddp.revenue_jpy * 0.04 / 12), 0),
    ]
    if ddp.jp_seed.owner_private_expense:
        rows.append(("役員貸付金", ddp.jp_seed.owner_loan_amount_jpy, 0))
        rows.append(("旅費交通費 (役員)", int(ddp.revenue_jpy * 0.018 / 12), 0))  # 高水準 = red flag
        rows.append(("接待交際費 (役員)", int(ddp.revenue_jpy * 0.020 / 12), 0))  # 高水準 = red flag

    for idx, (account, debit, credit) in enumerate(rows, start=5):
        ws.cell(row=idx, column=1, value=account)
        ws.cell(row=idx, column=2, value=debit)
        ws.cell(row=idx, column=3, value=credit)

    ws.column_dimensions["A"].width = 25
    ws.column_dimensions["B"].width = 18
    ws.column_dimensions["C"].width = 18
    wb.save(str(out_path))


def generate_xlsx_shareholder_register(ddp: DDProject, out_path: Path, rng: random.Random) -> None:
    """株主名簿 sample XLSX を literal 生成 (同族経営 / 名義株 pattern を visible 化)."""
    out_path.parent.mkdir(parents=True, exist_ok=True)
    wb = Workbook()
    ws = wb.active
    ws.title = "株主名簿"

    ws["A1"] = f"株主名簿 — {ddp.company_name}"
    ws["A1"].font = Font(bold=True, size=14)
    ws["A2"] = "基準日: 2025-03-31"
    headers = ["株主氏名", "所有株式数", "持分比率 (%)", "取得日", "備考"]
    for col_idx, h in enumerate(headers, start=1):
        cell = ws.cell(row=4, column=col_idx, value=h)
        cell.font = Font(bold=True)
        cell.fill = PatternFill("solid", fgColor="DDDDDD")

    total_shares = 100_000
    for idx, (name, ratio) in enumerate(ddp.shareholders, start=5):
        shares = int(total_shares * ratio / 100)
        ws.cell(row=idx, column=1, value=name)
        ws.cell(row=idx, column=2, value=shares)
        ws.cell(row=idx, column=3, value=ratio)
        ws.cell(row=idx, column=4, value=f"{rng.randint(2005, 2023)}-{rng.randint(1, 12):02d}-01")
        # 同姓家族 marker
        if ddp.jp_seed.family_governance and name.split(" ")[0] == ddp.jp_seed.family_surname:
            ws.cell(row=idx, column=5, value=f"創業者家族 ({ddp.jp_seed.family_surname}家)")

    # 名義変更履歴 sheet (名義株 pattern inject 時のみ)
    if ddp.jp_seed.nominee_shareholder and ddp.jp_seed.nominee_chain:
        ws2 = wb.create_sheet("名義変更履歴")
        ws2["A1"] = "名義変更履歴 (家族間)"
        ws2["A1"].font = Font(bold=True, size=14)
        ws2["A3"] = "日付"
        ws2["B3"] = "譲渡人"
        ws2["C3"] = "譲受人"
        ws2["D3"] = "備考"
        for col in ["A3", "B3", "C3", "D3"]:
            ws2[col].font = Font(bold=True)
        for idx, (frm, to, dt) in enumerate(ddp.jp_seed.nominee_chain, start=4):
            ws2.cell(row=idx, column=1, value=dt)
            ws2.cell(row=idx, column=2, value=frm)
            ws2.cell(row=idx, column=3, value=to)
            ws2.cell(row=idx, column=4, value=f"家族間譲渡 ({ddp.jp_seed.family_surname}家内)")
        ws2.column_dimensions["A"].width = 14
        ws2.column_dimensions["B"].width = 18
        ws2.column_dimensions["C"].width = 18
        ws2.column_dimensions["D"].width = 30

    ws.column_dimensions["A"].width = 22
    ws.column_dimensions["B"].width = 16
    ws.column_dimensions["C"].width = 14
    ws.column_dimensions["D"].width = 14
    ws.column_dimensions["E"].width = 30
    wb.save(str(out_path))


# ===== PPTX generators (python-pptx) =====


def generate_pptx_business_plan(ddp: DDProject, out_path: Path, rng: random.Random) -> None:
    """事業計画書 sample PPTX を literal 生成."""
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = f"中期事業計画 — {ddp.company_name}"
    if slide1.placeholders[1]:
        slide1.placeholders[1].text = f"2025-2027 年度 | {ddp.industry} | {ddp.location_pref}"

    # Slide 2: 現状業績
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "1. 現状業績"
    body = slide2.placeholders[1].text_frame
    body.text = f"売上高: ¥{ddp.revenue_jpy:,}"
    p = body.add_paragraph()
    p.text = f"従業員数: {ddp.employee_count} 名"
    p2 = body.add_paragraph()
    p2.text = f"主要事業地域: {ddp.location_pref}"
    if ddp.jp_seed.owner_private_expense:
        p3 = body.add_paragraph()
        median = industry_median_compensation_jpy(ddp.industry)
        p3.text = (
            f"※役員報酬: ¥{int(median * ddp.jp_seed.owner_compensation_multiplier):,}"
            f" (業界中央値 ¥{median:,} の {ddp.jp_seed.owner_compensation_multiplier} 倍)"
        )

    # Slide 3: 来期計画
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "2. 来期計画"
    body3 = slide3.placeholders[1].text_frame
    growth = rng.randint(5, 18)
    body3.text = f"売上目標: ¥{int(ddp.revenue_jpy * (1 + growth / 100)):,} (前年比 +{growth}%)"
    body3.add_paragraph().text = f"投資計画: ¥{int(ddp.revenue_jpy * rng.uniform(0.05, 0.12)):,}"
    body3.add_paragraph().text = f"新規採用: {rng.randint(5, 25)} 名"

    # Slide 4: 連絡先
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "3. 担当連絡先"
    body4 = slide4.placeholders[1].text_frame
    body4.text = f"代表取締役: {ddp.founder_name} ({ddp.founder_age} 歳)"
    body4.add_paragraph().text = f"Email: {ddp.contact_email}"
    body4.add_paragraph().text = f"Tel: {ddp.contact_phone}"

    prs.save(str(out_path))


def generate_pptx_shareholder_structure(ddp: DDProject, out_path: Path, rng: random.Random) -> None:
    """株主構成 sample PPTX を literal 生成 (同族経営 family tree visible 化)."""
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = f"株主構成 — {ddp.company_name}"
    if slide1.placeholders[1]:
        slide1.placeholders[1].text = f"基準日 2025-03-31"

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "1. 株主一覧"
    body = slide2.placeholders[1].text_frame
    first = True
    for name, ratio in ddp.shareholders[:8]:
        line = f"{name}: {ratio}%"
        if first:
            body.text = line
            first = False
        else:
            body.add_paragraph().text = line

    if ddp.jp_seed.family_governance:
        slide3 = prs.slides.add_slide(prs.slide_layouts[1])
        slide3.shapes.title.text = "2. 創業者家族の持分集中"
        body3 = slide3.placeholders[1].text_frame
        family_total = sum(
            ratio for name, ratio in ddp.shareholders if name.split(" ")[0] == ddp.jp_seed.family_surname
        )
        body3.text = f"{ddp.jp_seed.family_surname}家 合計持分: {family_total:.2f}%"
        body3.add_paragraph().text = "→ 同族支配下にあり、 経営判断は family 内合意で決定 (議事録参照)"
        body3.add_paragraph().text = f"主要家族構成員: {', '.join(ddp.jp_seed.family_members[:5])}"

    if ddp.jp_seed.nominee_shareholder:
        slide4 = prs.slides.add_slide(prs.slide_layouts[1])
        slide4.shapes.title.text = "3. 名義変更履歴 (家族間)"
        body4 = slide4.placeholders[1].text_frame
        first = True
        for frm, to, dt in ddp.jp_seed.nominee_chain:
            line = f"{dt}: {frm} → {to}"
            if first:
                body4.text = line
                first = False
            else:
                body4.add_paragraph().text = line

    prs.save(str(out_path))


# ===== orchestrator =====


def _generators() -> dict[str, list[tuple[str, Callable]]]:
    """doc_kind → [(filename, generator_fn), ...] の dict."""
    return {
        "pdf": [
            ("tax_return.pdf", generate_pdf_tax_return),
            ("financial_statement.pdf", generate_pdf_financial_statement),
        ],
        "docx": [
            ("board_minutes.docx", generate_docx_minutes),
            ("share_transfer_regulations.docx", generate_docx_regulations),
        ],
        "xlsx": [
            ("trial_balance.xlsx", generate_xlsx_trial_balance),
            ("shareholder_register.xlsx", generate_xlsx_shareholder_register),
        ],
        "pptx": [
            ("business_plan.pptx", generate_pptx_business_plan),
            ("shareholder_structure.pptx", generate_pptx_shareholder_structure),
        ],
    }


def generate_all_for_ddp(ddp: DDProject, output_root: Path, rng: random.Random) -> int:
    """1 DDP 用に全 doc 種別生成、 生成数を返却."""
    count = 0
    for kind, generators in _generators().items():
        kind_dir = output_root / ddp.ddp_id / kind
        for filename, fn in generators:
            out_path = kind_dir / filename
            fn(ddp, out_path, rng)
            count += 1
    return count


def main() -> int:
    rng = random.Random(SEED)
    faker = Faker("ja_JP")
    Faker.seed(SEED)

    OUTPUT_ROOT.mkdir(parents=True, exist_ok=True)
    summary: list[dict] = []
    print(f"[VDR-gen] SEED={SEED} / DDP_COUNT={DDP_COUNT} / OUTPUT_ROOT={OUTPUT_ROOT}")

    for i in range(DDP_COUNT):
        ddp = make_ddp(i, rng, faker)
        n = generate_all_for_ddp(ddp, OUTPUT_ROOT, rng)
        pattern_hits = []
        if ddp.jp_seed.family_governance:
            pattern_hits.append("family_governance")
        if ddp.jp_seed.nominee_shareholder:
            pattern_hits.append("nominee_shareholder")
        if ddp.jp_seed.owner_private_expense:
            pattern_hits.append("owner_private_expense")
        summary.append(
            {
                "ddp_id": ddp.ddp_id,
                "company_name": ddp.company_name,
                "industry": ddp.industry,
                "revenue_jpy": ddp.revenue_jpy,
                "docs_generated": n,
                "jp_patterns": pattern_hits,
            }
        )
        print(
            f"  [{ddp.ddp_id}] {ddp.company_name} ({ddp.industry}, ¥{ddp.revenue_jpy:,})"
            f" docs={n} patterns={pattern_hits}"
        )

    # summary を JSONL 出力 (Week 2+ で DDProject schema seed として再利用可)
    summary_path = OUTPUT_ROOT / "ddp_summary.jsonl"
    with open(summary_path, "w", encoding="utf-8") as f:
        import json as _json
        for s in summary:
            f.write(_json.dumps(s, ensure_ascii=False) + "\n")
    print(f"[VDR-gen] done、 total docs={sum(s['docs_generated'] for s in summary)}、 summary={summary_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
